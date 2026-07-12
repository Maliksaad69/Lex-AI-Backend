"""
Multi-format document parser.

Supports: PDF, DOCX, PPTX, XLSX, CSV, TXT, JSON, HTML, images (OCR),
          RTF, MSG (Outlook), and more via textract fallback.

Architecture:
  parse(file_path) -> ParsedDocument
    ├── text: cleaned full text
    ├── pages: list of per-page text (PDF only)
    ├── metadata: title, author, dates, language
    └── tables: extracted table data (CSV/XLSX/PDF tables)
"""

from __future__ import annotations

import csv
import io
import json
import re
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any, Optional

import pandas as pd
from PIL import Image


@dataclass
class ParsedDocument:
    """Unified output from any parser."""

    text: str = ""
    pages: list[str] = field(default_factory=list)
    metadata: dict[str, Any] = field(default_factory=dict)
    tables: list[list[list[str]]] = field(default_factory=list)  # [table[row[cell]]]
    file_path: str = ""
    file_type: str = ""
    file_size: int = 0
    error: Optional[str] = None


# ── Helpers ────────────────────────────────────────────────────────────


def _clean_text(text: str) -> str:
    """Normalise whitespace."""
    if not text:
        return ""
    text = text.replace("\r", "\n")
    text = re.sub(r"\n{3,}", "\n\n", text)
    text = re.sub(r"[ \t]+", " ", text)
    return text.strip()


def _file_type(path: str | Path) -> str:
    return Path(path).suffix.lower().lstrip(".")


# ── PDF Parser (pymupdf) ───────────────────────────────────────────────


def _parse_pdf(file_path: str, ocr_enabled: bool = True) -> ParsedDocument:
    """Extract text from PDF using pymupdf with OCR fallback."""
    import fitz

    doc = fitz.open(str(file_path))
    pages: list[str] = []
    tables: list[list[list[str]]] = []
    metadata: dict[str, Any] = {}

    # Extract PDF metadata
    py_meta = doc.metadata or {}
    metadata["title"] = py_meta.get("title", "")
    metadata["author"] = py_meta.get("author", "")
    metadata["subject"] = py_meta.get("subject", "")
    metadata["page_count"] = len(doc)

    for page_num, page in enumerate(doc):
        # 1. Direct text extraction
        text = _clean_text(page.get_text("text"))

        # 2. OCR fallback for image-only/scanned pages
        if len(text) < 30 and ocr_enabled:
            try:
                import pytesseract

                pix = page.get_pixmap(dpi=300)
                img = Image.open(io.BytesIO(pix.tobytes("png")))
                ocr_text = pytesseract.image_to_string(img)
                if ocr_text.strip():
                    text = _clean_text(ocr_text)
            except Exception:
                pass

        # 3. Extract tables from page
        try:
            page_tables = page.find_tables()
            if page_tables:
                for tab in page_tables:
                    tables.append(tab.extract())
        except Exception:
            pass

        pages.append(text)

    doc.close()

    full_text = "\n\n".join(pages)
    return ParsedDocument(
        text=full_text,
        pages=pages,
        metadata=metadata,
        tables=tables,
        file_path=str(file_path),
        file_type="pdf",
        file_size=Path(file_path).stat().st_size,
    )


# ── DOCX Parser ────────────────────────────────────────────────────────


def _parse_docx(file_path: str) -> ParsedDocument:
    """Extract text + tables from DOCX including headers/footers."""
    from docx import Document

    doc = Document(str(file_path))
    content: list[str] = []
    tables: list[list[list[str]]] = []
    metadata: dict[str, Any] = {}

    # Core properties
    props = doc.core_properties
    metadata["title"] = props.title or ""
    metadata["author"] = props.author or ""
    metadata["created"] = str(props.created) if props.created else ""

    # Paragraphs
    for para in doc.paragraphs:
        text = para.text.strip()
        if text:
            # Detect heading style for structure
            if para.style and para.style.name and para.style.name.startswith("Heading"):
                content.append(f"\n## {text}")
            else:
                content.append(text)

    # Tables
    for table in doc.tables:
        rows: list[list[str]] = []
        for row in table.rows:
            rows.append([cell.text.strip() for cell in row.cells])
        if rows:
            tables.append(rows)
        content.append("")  # separator

    return ParsedDocument(
        text="\n".join(content),
        metadata=metadata,
        tables=tables,
        file_path=str(file_path),
        file_type="docx",
        file_size=Path(file_path).stat().st_size,
    )


# ── PPTX Parser ────────────────────────────────────────────────────────


def _parse_pptx(file_path: str) -> ParsedDocument:
    """Extract text from PowerPoint slides."""
    from pptx import Presentation

    prs = Presentation(str(file_path))
    slides: list[str] = []
    metadata: dict[str, Any] = {"slide_count": len(prs.slides)}

    for i, slide in enumerate(prs.slides, 1):
        lines: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        lines.append(t)
            if shape.has_table:
                table = shape.table
                for row in table.rows:
                    row_text = " | ".join(cell.text.strip() for cell in row.cells)
                    lines.append(row_text)
        if lines:
            slides.append(f"--- Slide {i} ---\n" + "\n".join(lines))

    return ParsedDocument(
        text="\n\n".join(slides),
        metadata=metadata,
        file_path=str(file_path),
        file_type="pptx",
        file_size=Path(file_path).stat().st_size,
    )


# ── Spreadsheet Parsers ────────────────────────────────────────────────


def _parse_xlsx(file_path: str) -> ParsedDocument:
    """Extract all sheets from Excel as text + tables."""
    xl = pd.ExcelFile(str(file_path))
    sheets: list[str] = []
    tables: list[list[list[str]]] = []

    for sheet_name in xl.sheet_names:
        df = pd.read_excel(xl, sheet_name=sheet_name)
        sheets.append(f"--- Sheet: {sheet_name} ---\n{df.to_string(index=False)}")
        # Store as table: header + rows
        table = [df.columns.tolist()] + df.values.tolist()
        tables.append([[str(c) for c in row] for row in table])

    return ParsedDocument(
        text="\n\n".join(sheets),
        tables=tables,
        metadata={"sheets": xl.sheet_names},
        file_path=str(file_path),
        file_type="xlsx",
        file_size=Path(file_path).stat().st_size,
    )


def _parse_csv(file_path: str) -> ParsedDocument:
    """Parse CSV with auto-detected delimiter."""
    # Try reading with pandas first (auto-detect)
    try:
        df = pd.read_csv(str(file_path), nrows=None)
        text = df.to_string(index=False)
        table = [df.columns.tolist()] + df.values.tolist()
        tables = [[[str(c) for c in row] for row in table]]
    except Exception:
        # Fallback: read as plain text
        with open(file_path, encoding="utf-8", errors="replace") as f:
            text = f.read()
        tables = []

    return ParsedDocument(
        text=text,
        tables=tables,
        file_path=str(file_path),
        file_type="csv",
        file_size=Path(file_path).stat().st_size,
    )


# ── Image Parser (OCR only) ────────────────────────────────────────────


def _parse_image(file_path: str) -> ParsedDocument:
    """OCR an image file."""
    try:
        import pytesseract

        img = Image.open(str(file_path))
        text = pytesseract.image_to_string(img)
    except ImportError:
        text = ""
    except Exception as e:
        return ParsedDocument(
            text="",
            error=str(e),
            file_path=str(file_path),
            file_type=_file_type(file_path),
            file_size=Path(file_path).stat().st_size,
        )

    return ParsedDocument(
        text=_clean_text(text),
        file_path=str(file_path),
        file_type=_file_type(file_path),
        file_size=Path(file_path).stat().st_size,
    )


# ── Plain-text Parsers ─────────────────────────────────────────────────


def _parse_txt(file_path: str) -> ParsedDocument:
    with open(file_path, encoding="utf-8", errors="replace") as f:
        text = f.read()
    return ParsedDocument(
        text=text,
        file_path=str(file_path),
        file_type="txt",
        file_size=Path(file_path).stat().st_size,
    )


def _parse_json(file_path: str) -> ParsedDocument:
    with open(file_path, encoding="utf-8") as f:
        data = json.load(f)
    # Pretty-print for readable text
    text = json.dumps(data, indent=2, ensure_ascii=False)
    return ParsedDocument(
        text=text,
        file_path=str(file_path),
        file_type="json",
        file_size=Path(file_path).stat().st_size,
    )


def _parse_html(file_path: str) -> ParsedDocument:
    from bs4 import BeautifulSoup

    with open(file_path, encoding="utf-8", errors="replace") as f:
        soup = BeautifulSoup(f.read(), "lxml")
    # Remove script/style
    for tag in soup(["script", "style", "nav", "footer"]):
        tag.decompose()
    text = soup.get_text(separator="\n")
    return ParsedDocument(
        text=_clean_text(text),
        metadata={"title": soup.title.string if soup.title else ""},
        file_path=str(file_path),
        file_type="html",
        file_size=Path(file_path).stat().st_size,
    )


# ── Unified Parser ─────────────────────────────────────────────────────

PARSER_MAP = {
    "pdf": _parse_pdf,
    "docx": _parse_docx,
    "doc": _parse_docx,  # try python-docx for .doc too (may fail)
    "pptx": _parse_pptx,
    "ppt": _parse_pptx,
    "xlsx": _parse_xlsx,
    "xls": _parse_xlsx,
    "csv": _parse_csv,
    "txt": _parse_txt,
    "json": _parse_json,
    "html": _parse_html,
    "htm": _parse_html,
    "png": _parse_image,
    "jpg": _parse_image,
    "jpeg": _parse_image,
    "tiff": _parse_image,
    "bmp": _parse_image,
    "gif": _parse_image,
    "webp": _parse_image,
}


def parse(file_path: str | Path, **kwargs) -> ParsedDocument:
    """
    Parse any supported document into a ParsedDocument.

    Args:
        file_path: Path to the document.
        **kwargs: Passed to the specific parser (e.g., ocr_enabled=False).

    Returns:
        ParsedDocument with text, metadata, tables, pages.
    """
    file_path = Path(file_path)
    ftype = _file_type(file_path)

    if ftype not in PARSER_MAP:
        # Fallback: try textract for unsupported formats
        try:
            import textract

            raw = textract.process(str(file_path)).decode("utf-8", errors="replace")
            return ParsedDocument(
                text=_clean_text(raw),
                file_path=str(file_path),
                file_type=ftype,
                file_size=file_path.stat().st_size,
            )
        except Exception:
            return ParsedDocument(
                text="",
                error=f"Unsupported file type: .{ftype}",
                file_path=str(file_path),
                file_type=ftype,
                file_size=file_path.stat().st_size,
            )

    try:
        return PARSER_MAP[ftype](str(file_path), **kwargs)
    except Exception as e:
        # Last-resort fallback via textract
        try:
            import textract

            raw = textract.process(str(file_path)).decode("utf-8", errors="replace")
            return ParsedDocument(
                text=_clean_text(raw),
                file_path=str(file_path),
                file_type=ftype,
                file_size=file_path.stat().st_size,
            )
        except Exception:
            return ParsedDocument(
                text="",
                error=str(e),
                file_path=str(file_path),
                file_type=ftype,
                file_size=file_path.stat().st_size,
            )
